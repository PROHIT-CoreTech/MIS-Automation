"""
Downloads Page — Excel & PowerPoint MIS Report Export
Generates formatted reports from the live database and offers download buttons.
"""
import io
import sqlite3
import streamlit as st
import pandas as pd
from datetime import date, datetime
from core.auth import can_download_excel, can_download_ppt, is_admin

MONTHS = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun',
          'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']

DB_PATH = "data/mis_portal.db"

# ── Tally section mapping (mirrors reports.py) ─────────────
TALLY_SECTION = {
    'sales accounts':         'revenue',
    'direct incomes':         'dir_inc',
    'opening stock':          'opening',
    'purchase accounts':      'purchases',
    'add: purchase accounts': 'purchases',
    'direct expenses':        'direct_exp',
    'less: closing stock':    'closing',
    'closing stock':          'closing',
    'cost of sales :':        'cos_net',
    'indirect incomes':       'ind_inc',
    'indirect expenses':      'overhead',
    'salaries and bonus':     'overhead',
    'salary accounts':        'overhead',
}


def _get_conn():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn


def _load_pl_data(company_id: int, from_yr: int, from_mo: int,
                  to_yr: int, to_mo: int):
    """Return P&L rows and selected month labels."""
    conn = _get_conn()
    avail = conn.execute(
        "SELECT DISTINCT year, month FROM pl_data "
        "WHERE company_id=? ORDER BY year, month", (company_id,)
    ).fetchall()

    sel_months = [
        (int(r['year']), int(r['month'])) for r in avail
        if (from_yr, from_mo) <= (int(r['year']), int(r['month'])) <= (to_yr, to_mo)
    ]
    mo_labels = [f"{MONTHS[m-1]}-{str(y)[2:]}" for y, m in sel_months]

    rows = conn.execute("""
        SELECT ledger_name, tally_group, mis_group, year, month, net
        FROM pl_data
        WHERE company_id=?
          AND ((year > ?) OR (year = ? AND month >= ?))
          AND ((year < ?) OR (year = ? AND month <= ?))
        ORDER BY tally_group, ledger_name, year, month
    """, (company_id, from_yr, from_yr, from_mo,
          to_yr, to_yr, to_mo)).fetchall()
    conn.close()
    return rows, mo_labels, sel_months


def _build_sections(rows, mo_labels):
    """Group P&L rows into section → ledger → {month: value} dict."""
    sections = {}
    for r in rows:
        tg  = (r['tally_group'] or '').lower().strip()
        sec = TALLY_SECTION.get(tg, 'other')
        lbl = f"{MONTHS[int(r['month'])-1]}-{str(int(r['year']))[2:]}"
        if lbl not in mo_labels:
            continue
        ln  = r['ledger_name']
        sections.setdefault(sec, {}).setdefault(ln, {})[lbl] = float(r['net'] or 0)
    return sections


def _fmt_cr(v: float) -> str:
    if v == 0:
        return "—"
    sign = "-" if v < 0 else ""
    a = abs(v)
    if a >= 1e7:
        return f"{sign}₹{a/1e7:.2f} Cr"
    elif a >= 1e5:
        return f"{sign}₹{a/1e5:.2f} L"
    elif a >= 1000:
        return f"{sign}₹{a/1000:.1f}K"
    return f"{sign}₹{a:,.0f}"


# ── EXCEL EXPORT ───────────────────────────────────────────
def _generate_excel(company_name: str, from_lbl: str, to_lbl: str,
                    sections: dict, mo_labels: list) -> bytes:
    """Generate a formatted .xlsx P&L report using openpyxl."""
    from openpyxl import Workbook
    from openpyxl.styles import (Font, PatternFill, Alignment,
                                  Border, Side)
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    ws = wb.active
    ws.title = "P&L Report"

    def fill(hex_):
        return PatternFill('solid', start_color=hex_, end_color=hex_)
    def thin_border():
        s = Side(style='thin', color='CCCCCC')
        return Border(left=s, right=s, top=s, bottom=s)

    H_FONT  = Font(name='Arial', bold=True, color='FFFFFF', size=10)
    S_FONT  = Font(name='Arial', bold=True, color='FFFFFF', size=10)
    ST_FONT = Font(name='Arial', bold=True, color='1D3557', size=9)
    L_FONT  = Font(name='Arial', size=9)
    GP_FONT = Font(name='Arial', bold=True, color='FFFFFF', size=10)
    NP_FONT = Font(name='Arial', bold=True, color='FFFFFF', size=10)
    NUM_FMT = '#,##0.00'

    H_FILL  = fill('1D3557')
    S_FILL  = fill('2E75B6')
    ST_FILL = fill('D6E4F0')
    GP_FILL = fill('1D9E75')
    NP_FILL = fill('D85A30')
    L_FILL  = fill('FFFFFF')
    A_FILL  = fill('F8FAFB')

    # Title
    n_cols = len(mo_labels) + 2
    ws.merge_cells(f'A1:{get_column_letter(n_cols)}1')
    ws['A1'] = f"P&L Report | {company_name} | {from_lbl} → {to_lbl}"
    ws['A1'].font      = Font(name='Arial', bold=True, size=13, color='1D3557')
    ws['A1'].alignment = Alignment(horizontal='center', vertical='center')
    ws.row_dimensions[1].height = 28

    # Column headers
    headers = ['Particulars'] + mo_labels + ['Total']
    for ci, h in enumerate(headers, 1):
        c = ws.cell(row=2, column=ci, value=h)
        c.font      = H_FONT
        c.fill      = H_FILL
        c.alignment = Alignment(horizontal='left' if ci == 1 else 'right',
                                vertical='center')
        c.border    = thin_border()
    ws.row_dimensions[2].height = 20
    ws.column_dimensions['A'].width = 38
    for ci in range(2, n_cols + 1):
        ws.column_dimensions[get_column_letter(ci)].width = 13
    ws.freeze_panes = 'B3'

    row_idx = [3]

    def write_row(label, values, font, bg_fill, indent=0):
        r = row_idx[0]
        prefix = '    ' * indent
        c0 = ws.cell(row=r, column=1, value=prefix + label)
        c0.font = font; c0.fill = bg_fill
        c0.alignment = Alignment(horizontal='left', vertical='center')
        c0.border = thin_border()
        for ci, v in enumerate(values, 2):
            c = ws.cell(row=r, column=ci, value=v if v else None)
            c.font = font; c.fill = bg_fill
            c.alignment = Alignment(horizontal='right', vertical='center')
            c.border = thin_border()
            if isinstance(v, (int, float)) and v:
                c.number_format = NUM_FMT
        ws.row_dimensions[r].height = 16
        row_idx[0] += 1

    def sec_vals(sec_key):
        data = sections.get(sec_key, {})
        totals = {}
        for ln, mo_data in data.items():
            for lbl, v in mo_data.items():
                totals[lbl] = totals.get(lbl, 0) + abs(v)
        return [round(totals.get(l, 0) / 1e7, 4) for l in mo_labels]

    def add_section(title, sec_key):
        data = sections.get(sec_key, {})
        if not data:
            return
        write_row(title, [''] * (len(mo_labels) + 1), S_FONT, S_FILL)
        alt = False
        for ln, mo_data in sorted(data.items()):
            vals = [round(abs(mo_data.get(l, 0)) / 1e7, 4)
                    if abs(mo_data.get(l, 0)) >= 100 else None for l in mo_labels]
            total = sum(abs(mo_data.get(l, 0)) for l in mo_labels)
            vals.append(round(total / 1e7, 4) if total >= 100 else None)
            write_row(ln, vals, L_FONT, A_FILL if alt else L_FILL, indent=1)
            alt = not alt
        sub = sec_vals(sec_key)
        sub.append(round(sum(sub), 4))
        write_row('Sub-Total', sub, ST_FONT, ST_FILL)

    add_section('SALES ACCOUNTS (Revenue)', 'revenue')
    add_section('DIRECT INCOMES', 'dir_inc')

    # COGS
    write_row('COST OF GOODS SOLD', [''] * (len(mo_labels) + 1), S_FONT, S_FILL)
    for sub_title, sub_key in [
        ('Opening Stock',          'opening'),
        ('Add: Purchase Accounts', 'purchases'),
        ('Add: Direct Expenses',   'direct_exp'),
        ('Less: Closing Stock',    'closing'),
    ]:
        if sections.get(sub_key):
            add_section(sub_title, sub_key)

    # GP
    gp_mo = [round((sec_vals('revenue')[i] + sec_vals('dir_inc')[i]
                    - sec_vals('opening')[i] - sec_vals('purchases')[i]
                    - sec_vals('direct_exp')[i] + sec_vals('closing')[i]), 4)
             for i in range(len(mo_labels))]
    gp_mo.append(round(sum(gp_mo), 4))
    write_row('GROSS PROFIT (c/o)', gp_mo, GP_FONT, GP_FILL)

    add_section('INDIRECT INCOMES', 'ind_inc')
    add_section('INDIRECT EXPENSES (Overhead)', 'overhead')

    # NP
    np_mo = [round((sec_vals('revenue')[i] + sec_vals('dir_inc')[i]
                    - sec_vals('opening')[i] - sec_vals('purchases')[i]
                    - sec_vals('direct_exp')[i] + sec_vals('closing')[i]
                    + sec_vals('ind_inc')[i] - sec_vals('overhead')[i]), 4)
             for i in range(len(mo_labels))]
    np_mo.append(round(sum(np_mo), 4))
    write_row('NET PROFIT (Nett)', np_mo, NP_FONT, NP_FILL)

    # Note
    row_idx[0] += 1
    ws.cell(row=row_idx[0], column=1,
            value="Note: All values in Crores (Cr). Generated by MIS Portal.")
    ws.cell(row=row_idx[0], column=1).font = Font(
        name='Arial', size=8, italic=True, color='888888')

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.getvalue()


# ── PPT EXPORT ─────────────────────────────────────────────
def _generate_ppt(company_name: str, from_lbl: str, to_lbl: str,
                  sections: dict, mo_labels: list) -> bytes:
    """Generate a structured .pptx MIS summary report using python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    DARK   = RGBColor(0x1D, 0x35, 0x57)
    ACCENT = RGBColor(0x2E, 0x75, 0xB6)
    GREEN  = RGBColor(0x1D, 0x9E, 0x75)
    RED    = RGBColor(0xD8, 0x5A, 0x30)
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    LGRAY  = RGBColor(0xF0, 0xF4, 0xF8)

    blank_layout = prs.slide_layouts[6]  # completely blank

    def add_textbox(slide, text, left, top, width, height,
                    font_size=14, bold=False, color=None,
                    bg_color=None, align=PP_ALIGN.LEFT, italic=False):
        txb = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height))
        tf  = txb.text_frame
        tf.word_wrap = True
        p   = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(font_size)
        run.font.bold  = bold
        run.font.italic = italic
        run.font.color.rgb = color or DARK
        if bg_color:
            from pptx.util import Pt as _Pt
            fill = txb.fill
            fill.solid()
            fill.fore_color.rgb = bg_color
        return txb

    def add_rect(slide, left, top, width, height, fill_color, line_color=None):
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if line_color:
            shape.line.color.rgb = line_color
        else:
            shape.line.fill.background()
        return shape

    def sec_total(sec_key):
        data = sections.get(sec_key, {})
        total = 0.0
        for mo_data in data.values():
            total += sum(abs(v) for v in mo_data.values())
        return total / 1e7

    # ── SLIDE 1: Title ─────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, 13.33, 7.5, DARK)
    add_rect(slide, 0, 5.8, 13.33, 1.7, ACCENT)
    add_textbox(slide, "MIS REPORT", 1, 1.2, 11, 1.2,
                font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, company_name, 1, 2.8, 11, 0.9,
                font_size=24, bold=False, color=RGBColor(0xA0, 0xC4, 0xFF),
                align=PP_ALIGN.CENTER)
    add_textbox(slide, f"Period: {from_lbl}  →  {to_lbl}", 1, 3.7, 11, 0.6,
                font_size=16, color=RGBColor(0xCC, 0xDD, 0xFF),
                align=PP_ALIGN.CENTER)
    add_textbox(slide, f"Generated on {datetime.now().strftime('%d %b %Y')}",
                1, 6.0, 11, 0.5, font_size=12,
                color=WHITE, align=PP_ALIGN.CENTER, italic=True)

    # ── SLIDE 2: P&L Summary KPIs ──────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, 13.33, 1.0, DARK)
    add_textbox(slide, f"P&L Summary — {company_name}  |  {from_lbl} → {to_lbl}",
                0.2, 0.1, 12, 0.8, font_size=16, bold=True, color=WHITE)

    kpis = [
        ("Revenue",       sec_total('revenue'), ACCENT),
        ("Gross Profit",  sec_total('revenue') + sec_total('dir_inc')
                          - sec_total('opening') - sec_total('purchases')
                          - sec_total('direct_exp') + sec_total('closing'), GREEN),
        ("Indirect Inc.", sec_total('ind_inc'), RGBColor(0x00, 0xA3, 0xFF)),
        ("Overhead",      sec_total('overhead'), RGBColor(0xFF, 0xB3, 0x47)),
        ("Net Profit",    sec_total('revenue') + sec_total('dir_inc')
                          - sec_total('opening') - sec_total('purchases')
                          - sec_total('direct_exp') + sec_total('closing')
                          + sec_total('ind_inc') - sec_total('overhead'), RED),
    ]

    card_w = 2.4
    gap    = 0.13
    for i, (label, value, color) in enumerate(kpis):
        x = 0.15 + i * (card_w + gap)
        add_rect(slide, x, 1.3, card_w, 2.0, color)
        add_textbox(slide, label,
                    x + 0.1, 1.4, card_w - 0.2, 0.5,
                    font_size=11, bold=True, color=WHITE)
        sign = "-" if value < 0 else ""
        add_textbox(slide, f"{sign}₹{abs(value):.2f} Cr",
                    x + 0.1, 1.95, card_w - 0.2, 0.8,
                    font_size=20, bold=True, color=WHITE)

    # Monthly trend table
    add_rect(slide, 0.15, 3.6, 13.0, 0.4, ACCENT)
    add_textbox(slide, "Monthly Revenue (₹ Cr)", 0.25, 3.65, 5, 0.3,
                font_size=11, bold=True, color=WHITE)

    rev_data = sections.get('revenue', {})
    y_pos = 4.1
    for i, lbl in enumerate(mo_labels[:12]):  # max 12 months
        col_rev = sum(abs(d.get(lbl, 0)) for d in rev_data.values()) / 1e7
        x = 0.15 + i * (13.0 / max(len(mo_labels[:12]), 1))
        w = 13.0 / max(len(mo_labels[:12]), 1) - 0.05
        add_rect(slide, x, y_pos, w, 0.35,
                 RGBColor(0xE8, 0xF4, 0xF8))
        add_textbox(slide, lbl,
                    x, y_pos, w, 0.2, font_size=7, color=DARK,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, f"{col_rev:.2f}",
                    x, y_pos + 0.18, w, 0.2, font_size=8, bold=True,
                    color=ACCENT, align=PP_ALIGN.CENTER)

    # ── SLIDE 3: Detailed P&L Table ────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, 13.33, 1.0, DARK)
    add_textbox(slide, "Detailed P&L Breakup (₹ Crores)",
                0.2, 0.1, 12, 0.8, font_size=16, bold=True, color=WHITE)

    section_rows = [
        ("SALES ACCOUNTS",      sec_total('revenue'),   ACCENT),
        ("DIRECT INCOMES",      sec_total('dir_inc'),   ACCENT),
        ("COST OF GOODS SOLD",  sec_total('opening') + sec_total('purchases') +
                                sec_total('direct_exp') - sec_total('closing'), RED),
        ("GROSS PROFIT",        sec_total('revenue') + sec_total('dir_inc')
                                - sec_total('opening') - sec_total('purchases')
                                - sec_total('direct_exp') + sec_total('closing'), GREEN),
        ("INDIRECT INCOMES",    sec_total('ind_inc'),   ACCENT),
        ("INDIRECT EXPENSES",   sec_total('overhead'),  RGBColor(0xFF, 0xB3, 0x47)),
        ("NET PROFIT",          sec_total('revenue') + sec_total('dir_inc')
                                - sec_total('opening') - sec_total('purchases')
                                - sec_total('direct_exp') + sec_total('closing')
                                + sec_total('ind_inc') - sec_total('overhead'), RED),
    ]

    y = 1.2
    for label, value, color in section_rows:
        add_rect(slide, 0.3, y, 8.5, 0.55, LGRAY,
                 line_color=RGBColor(0xCC, 0xCC, 0xCC))
        add_textbox(slide, label, 0.4, y + 0.08, 6, 0.4,
                    font_size=10, bold=True, color=DARK)
        sign = "-" if value < 0 else ""
        add_textbox(slide, f"{sign}₹{abs(value):.2f} Cr",
                    6.5, y + 0.08, 2.2, 0.4,
                    font_size=10, bold=True, color=color,
                    align=PP_ALIGN.RIGHT)
        y += 0.62

    add_textbox(slide,
                f"Period: {from_lbl} → {to_lbl}  |  All values in Crores",
                0.3, 7.1, 12, 0.3, font_size=8,
                color=RGBColor(0x88, 0x88, 0x88), italic=True)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


# ── MAIN PAGE ──────────────────────────────────────────────
def show_downloads(user):
    st.markdown("## 📥 Downloads")
    st.caption("Export MIS reports in Excel or PowerPoint format.")

    # ── Read from shared sidebar filter (set by app.py) ──────
    company_id   = st.session_state.get('global_company_id')
    company_name = st.session_state.get('global_company_name', '')
    from_lbl     = st.session_state.get('global_from_lbl')
    to_lbl       = st.session_state.get('global_to_lbl')
    from_yr      = st.session_state.get('global_from_yr')
    from_mo      = st.session_state.get('global_from_mo')
    to_yr        = st.session_state.get('global_to_yr')
    to_mo        = st.session_state.get('global_to_mo')

    if not company_id or not from_lbl:
        st.info("Select a company and date range from the sidebar to generate reports.")
        return

    if (from_yr, from_mo) > (to_yr, to_mo):
        st.error("'From' cannot be after 'To'.")
        return

    # Summary header showing what's selected
    st.markdown(f"""
        <div style='background:rgba(99,102,241,0.08);border:1px solid rgba(99,102,241,0.2);
             border-radius:12px;padding:10px 16px;margin-bottom:1rem;
             display:flex;align-items:center;gap:16px;flex-wrap:wrap;'>
            <span style='font-weight:700;color:#a5b4fc;'>🏢 {company_name}</span>
            <span style='color:#475569;'>|</span>
            <span style='color:#94a3b8;font-size:0.88rem;'>📅 {from_lbl} → {to_lbl}</span>
        </div>
    """, unsafe_allow_html=True)

    st.markdown("---")

    # ── Build data ─────────────────────────────────────────
    with st.spinner("Preparing report data..."):
        rows, mo_labels, _ = _load_pl_data(
            company_id, from_yr, from_mo, to_yr, to_mo)
        sections = _build_sections(rows, mo_labels)

    n_months = len(mo_labels)
    st.caption(f"📊 {n_months} months selected: **{from_lbl}** → **{to_lbl}**")

    # ── Download cards ─────────────────────────────────────
    c1, c2 = st.columns(2)

    with c1:
        st.markdown("### 📗 Excel Export")
        st.markdown(
            "Full **P&L Detailed Report** with monthly columns, "
            "section headers, subtotals, GP and NP rows — formatted with colors."
        )
        if can_download_excel(user):
            if st.button("⚙️ Generate Excel", key="gen_xlsx",
                         use_container_width=True, type="primary"):
                with st.spinner("Building Excel..."):
                    try:
                        xls_bytes = _generate_excel(
                            company_name, from_lbl, to_lbl, sections, mo_labels)
                        fname = f"MIS_PL_{company_name}_{from_lbl}_to_{to_lbl}.xlsx"
                        for ch in '\\/:*?"<>| ':
                            fname = fname.replace(ch, '_')
                        st.download_button(
                            label="⬇️ Download Excel (.xlsx)",
                            data=xls_bytes,
                            file_name=fname,
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                            use_container_width=True,
                            key="dl_xlsx_btn",
                        )
                        st.success("✅ Excel ready!")
                    except Exception as e:
                        st.error(f"Excel generation failed: {e}")
        else:
            st.info("🔒 Excel download not enabled for your account.")

    with c2:
        st.markdown("### 📑 PowerPoint Export")
        st.markdown(
            "3-slide **MIS Presentation** with title slide, KPI summary cards, "
            "and detailed P&L table — ready to share with management."
        )
        if can_download_ppt(user):
            if st.button("⚙️ Generate PowerPoint", key="gen_pptx",
                         use_container_width=True, type="primary"):
                with st.spinner("Building PowerPoint..."):
                    try:
                        ppt_bytes = _generate_ppt(
                            company_name, from_lbl, to_lbl, sections, mo_labels)
                        fname = f"MIS_Report_{company_name}_{from_lbl}_to_{to_lbl}.pptx"
                        for ch in '\\/:*?"<>| ':
                            fname = fname.replace(ch, '_')
                        st.download_button(
                            label="⬇️ Download PowerPoint (.pptx)",
                            data=ppt_bytes,
                            file_name=fname,
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True,
                            key="dl_pptx_btn",
                        )
                        st.success("✅ PowerPoint ready!")
                    except Exception as e:
                        st.error(f"PowerPoint generation failed: {e}")
        else:
            st.info("🔒 PowerPoint download not enabled for your account.")
